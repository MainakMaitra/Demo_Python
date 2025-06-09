import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import datetime

def create_comprehensive_insights_presentation():
    """Create a comprehensive PowerPoint presentation with all insights from Insights_V2.md"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    DARK_BLUE = RGBColor(30, 60, 114)
    MEDIUM_BLUE = RGBColor(42, 82, 152)
    LIGHT_BLUE = RGBColor(232, 244, 248)
    RED = RGBColor(255, 107, 107)
    ORANGE = RGBColor(255, 167, 38)
    GREEN = RGBColor(102, 187, 106)
    GRAY = RGBColor(102, 102, 102)
    DARK_GRAY = RGBColor(51, 51, 51)
    YELLOW_BG = RGBColor(255, 243, 205)
    
    def add_slide_number(slide, number):
        """Add slide number to top right corner"""
        num_box = slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(0.8), Inches(0.3))
        num_frame = num_box.text_frame
        num_frame.text = str(number)
        num_frame.paragraphs[0].font.size = Pt(12)
        num_frame.paragraphs[0].font.color.rgb = GRAY
        num_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Holistic Precision Drop Analysis"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Critical Insights & Strategic Action Plan"
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    date_frame = date_box.text_frame
    date_frame.text = f"Executive Presentation - {datetime.datetime.now().strftime('%B %Y')}"
    date_frame.paragraphs[0].font.size = Pt(18)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary with All Key Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 2)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Current state box
    state_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.8)
    )
    state_box.fill.solid()
    state_box.fill.fore_color.rgb = LIGHT_BLUE
    
    state_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6))
    state_frame = state_text.text_frame
    state_frame.text = "Current State: 74.8% overall precision (4.8% above 70% target but declining)"
    state_frame.paragraphs[0].font.size = Pt(16)
    state_frame.paragraphs[0].font.bold = True
    
    # Key metrics
    metrics = [
        ("96.7%", "FPs from Context Issues", "701/725 records", RED),
        ("$2.4M", "Annual Cost Impact", "$100 per false escalation", ORANGE),
        ("37/158", "Categories Below Target", "23.4% failing", RED),
        ("83.6%", "Validation Agreement", "1.4% below target", ORANGE)
    ]
    
    x_positions = [0.5, 2.75, 5, 7.25]
    for i, (value, label, detail, color) in enumerate(metrics):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(2.2),
            Inches(2), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        value_box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(2.4),
            Inches(2), Inches(0.6)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(32)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = color
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(3),
            Inches(2), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(11)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        detail_box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(3.4),
            Inches(2), Inches(0.3)
        )
        detail_frame = detail_box.text_frame
        detail_frame.text = detail
        detail_frame.paragraphs[0].font.size = Pt(10)
        detail_frame.paragraphs[0].font.color.rgb = GRAY
        detail_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Strategic priority
    priority_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.8),
        Inches(9), Inches(1.2)
    )
    priority_box.fill.solid()
    priority_box.fill.fore_color.rgb = YELLOW_BG
    priority_box.line.color.rgb = RED
    priority_box.line.width = Pt(3)
    
    priority_text = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(8.4), Inches(1))
    priority_frame = priority_text.text_frame
    priority_frame.text = "Critical Finding: 96.7% of False Positives stem from context-insensitive rule processing"
    priority_frame.paragraphs[0].font.size = Pt(18)
    priority_frame.paragraphs[0].font.bold = True
    p = priority_frame.add_paragraph()
    p.text = "Business Impact: $2.4M annual cost with systematic erosion across 37/158 categories"
    p.font.size = Pt(14)
    p = priority_frame.add_paragraph()
    p.text = "Strategic Priority: Immediate technical intervention combined with operational standardization"
    p.font.size = Pt(14)
    p.font.italic = True
    
    # Slide 3: Performance vs Targets Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 3)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Key Performance vs Targets"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Performance table
    rows = 4
    cols = 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(2)).table
    
    # Headers
    headers = ['KPI Level', 'Metric', 'Target', 'Current', 'Status']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Data with status colors
    kpi_data = [
        ['Primary', 'Overall precision', '≥ 70%', '74.8%', 'EXCEEDING TARGET'],
        ['Secondary', 'All categories precision', '≥ 60%', '37/158 below 70%', 'FAILING'],
        ['Tertiary', 'Validation agreement', '≥ 85%', '83.6%', 'SLIGHTLY BELOW']
    ]
    
    status_colors = {
        'EXCEEDING TARGET': GREEN,
        'FAILING': RED,
        'SLIGHTLY BELOW': ORANGE
    }
    
    for i, row_data in enumerate(kpi_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = value
            if j == 4:  # Status column
                cell.fill.solid()
                cell.fill.fore_color.rgb = status_colors.get(value, GRAY)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Monthly precision trend
    chart_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.8),
        Inches(9), Inches(3)
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Add trend data
    trend_text = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(8.4), Inches(2.6))
    trend_frame = trend_text.text_frame
    trend_frame.text = "Monthly Precision Trend:"
    trend_frame.paragraphs[0].font.bold = True
    trend_frame.paragraphs[0].font.size = Pt(14)
    
    # Monthly data points
    monthly_data = [
        "Oct 2024: 80.1% ↗",
        "Nov 2024: 80.3% → (stable)",
        "Dec 2024: 70.5% ↘ (-9.8% MoM drop)",
        "Jan 2025: 75.2% ↗ (recovery)",
        "Feb 2025: 70.8% ↘ (-4.5% decline)",
        "Mar 2025: 71.9% ↗ (slight improvement)"
    ]
    
    for i, data_point in enumerate(monthly_data):
        p = trend_frame.add_paragraph()
        p.text = f"    • {data_point}"
        p.font.size = Pt(12)
        if "drop" in data_point or "decline" in data_point:
            p.font.color.rgb = RED
        elif "recovery" in data_point or "improvement" in data_point:
            p.font.color.rgb = GREEN
    
    # Slide 4: Context-Insensitive Negation (Critical Priority #1)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 4)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Critical Priority #1: Context-Insensitive Negation"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Impact box
    impact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.8)
    )
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    impact_box.line.color.rgb = RED
    impact_box.line.width = Pt(3)
    
    impact_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6))
    impact_frame = impact_text.text_frame
    impact_frame.text = "CRITICAL: 96.7% of all FPs (701/725 records) | Financial Impact: $2.4M annual cost"
    impact_frame.paragraphs[0].font.size = Pt(16)
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.color.rgb = RED
    
    # Evidence data
    evidence_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2.2),
        Inches(4.2), Inches(2)
    )
    evidence_box.fill.solid()
    evidence_box.fill.fore_color.rgb = RGBColor(255, 250, 250)
    
    evidence_text = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(3.6), Inches(1.6))
    evidence_frame = evidence_text.text_frame
    evidence_frame.text = "Negation Pattern Analysis:"
    evidence_frame.paragraphs[0].font.bold = True
    evidence_frame.paragraphs[0].font.size = Pt(14)
    
    evidence_data = [
        "True Positives: 11.789 negations/transcript",
        "False Positives: 6.233 negations/transcript",
        "Risk Factor: 0.529",
        "Interpretation: FPs have fewer negations"
    ]
    
    for item in evidence_data:
        p = evidence_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        if "Risk Factor" in item:
            p.font.bold = True
    
    # Root cause box
    cause_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(2.2),
        Inches(4.5), Inches(2)
    )
    cause_box.fill.solid()
    cause_box.fill.fore_color.rgb = RGBColor(255, 253, 240)
    
    cause_text = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.9), Inches(1.6))
    cause_frame = cause_text.text_frame
    cause_frame.text = "Root Cause:"
    cause_frame.paragraphs[0].font.bold = True
    cause_frame.paragraphs[0].font.size = Pt(14)
    p = cause_frame.add_paragraph()
    p.text = "Rules trigger on complaint keywords without understanding negation context"
    p.font.size = Pt(12)
    p = cause_frame.add_paragraph()
    p.text = "\nExample Pattern:"
    p.font.bold = True
    p.font.size = Pt(12)
    p = cause_frame.add_paragraph()
    p.text = '"I\'m NOT complaining, but..."'
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    
    # Solution implementation
    solution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.5),
        Inches(9), Inches(2)
    )
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = LIGHT_BLUE
    
    solution_text = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(8.4), Inches(1.6))
    solution_frame = solution_text.text_frame
    solution_frame.text = "Immediate Actions:"
    solution_frame.paragraphs[0].font.bold = True
    solution_frame.paragraphs[0].font.size = Pt(14)
    
    solutions = [
        "1. Universal negation template: (query) AND NOT ((not|no|never|don't) NEAR:3 (complain|complaint))",
        "2. Context window expansion: (complaint_terms) NOT WITHIN:10 (negation_patterns)",
        "3. Deploy to top 5 worst-performing categories first"
    ]
    
    for solution in solutions:
        p = solution_frame.add_paragraph()
        p.text = solution
        p.font.size = Pt(11)
        if "Universal negation" in solution or "Context window" in solution:
            p.font.name = 'Courier New'
    
    # Expected impact
    impact_data = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.4))
    impact_frame = impact_data.text_frame
    impact_frame.text = "Expected Impact: +15-20% precision | Effort: Medium | Timeline: 2-4 weeks"
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.size = Pt(13)
    impact_frame.paragraphs[0].font.color.rgb = GREEN
    
    # Slide 5: Agent Explanation Contamination (Critical Priority #2)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 5)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Critical Priority #2: Agent Explanation Contamination"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Impact summary
    impact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.8)
    )
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
    impact_box.line.color.rgb = ORANGE
    impact_box.line.width = Pt(3)
    
    impact_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6))
    impact_frame = impact_text.text_frame
    impact_frame.text = "HIGH PRIORITY: 58.5% of all FPs (424/725 records)"
    impact_frame.paragraphs[0].font.size = Pt(16)
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.color.rgb = ORANGE
    
    # Contamination rates table
    rows = 5
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.2), Inches(5), Inches(2.2)).table
    
    # Headers
    headers = ['Category', 'Contamination Rate', 'Severity']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    
    contamination_data = [
        ['Fees & Interest', '75.0%', 'Highest'],
        ['Billing Disputes', '66.2%', 'High'],
        ['Customer Relations', '54.9%', 'Medium'],
        ['Overall Average', '58.5%', 'High']
    ]
    
    severity_colors = {
        'Highest': RED,
        'High': ORANGE,
        'Medium': YELLOW_BG
    }
    
    for i, row_data in enumerate(contamination_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = value
            if j == 2:  # Severity column
                cell.fill.solid()
                cell.fill.fore_color.rgb = severity_colors.get(value, GRAY)
                if value in ['Highest', 'High']:
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Root cause explanation
    cause_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.8), Inches(2.2),
        Inches(3.7), Inches(2.2)
    )
    cause_box.fill.solid()
    cause_box.fill.fore_color.rgb = RGBColor(255, 250, 245)
    
    cause_text = slide.shapes.add_textbox(Inches(6), Inches(2.4), Inches(3.3), Inches(1.8))
    cause_frame = cause_text.text_frame
    cause_frame.text = "Root Cause:"
    cause_frame.paragraphs[0].font.bold = True
    cause_frame.paragraphs[0].font.size = Pt(13)
    p = cause_frame.add_paragraph()
    p.text = "Agent hypothetical scenarios trigger complaint detection"
    p.font.size = Pt(11)
    p = cause_frame.add_paragraph()
    p.text = "\nExample:"
    p.font.bold = True
    p.font.size = Pt(11)
    p = cause_frame.add_paragraph()
    p.text = '"If you were to complain about fees..."'
    p.font.name = 'Courier New'
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    
    # Solutions
    solution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.7),
        Inches(9), Inches(1.8)
    )
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = LIGHT_BLUE
    
    solution_text = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(8.4), Inches(1.4))
    solution_frame = solution_text.text_frame
    solution_frame.text = "Immediate Actions:"
    solution_frame.paragraphs[0].font.bold = True
    solution_frame.paragraphs[0].font.size = Pt(14)
    
    actions = [
        "1. Implement speaker role detection",
        "2. Add filters: AND NOT ((explain|example|suppose|hypothetically) NEAR:5 (complaint))",
        "3. Channel-specific optimization (focus on 'customer' channel only)"
    ]
    
    for action in actions:
        p = solution_frame.add_paragraph()
        p.text = action
        p.font.size = Pt(11)
        if "Add filters" in action:
            p.font.name = 'Courier New'
    
    # Expected outcome
    outcome_text = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.4))
    outcome_frame = outcome_text.text_frame
    outcome_frame.text = "Expected Impact: +8-12% precision | Effort: Low | Timeline: 1-2 weeks | Success Metric: <30% contamination"
    outcome_frame.paragraphs[0].font.bold = True
    outcome_frame.paragraphs[0].font.size = Pt(13)
    outcome_frame.paragraphs[0].font.color.rgb = GREEN
    
    # Slide 6: Emergency Category Review
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 6)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Emergency Category Review - URGENT"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RED
    
    # Critical categories with MoM drops
    rows = 4
    cols = 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(1.8)).table
    
    headers = ['Category', 'Precision', 'MoM Drop', 'Impact Score', 'Action']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    
    critical_data = [
        ['Customer Relations - Close Account', '56.2%', '-38.6%', '22.4', 'Add context filters'],
        ['Credit Card ATM - Unclassified', '39.6%', 'N/A', '14.6', 'Improve classification'],
        ['Credit Card ATM - Rejected/Declined', '47.3%', '-88.2%', '12.5', 'Separate technical issues']
    ]
    
    for i, row_data in enumerate(critical_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = value
            if j == 1 and float(value.strip('%')) < 50:  # Precision column
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 235, 235)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Specific findings
    findings_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.2),
        Inches(9), Inches(3.3)
    )
    findings_box.fill.solid()
    findings_box.fill.fore_color.rgb = RGBColor(255, 250, 250)
    
    findings_text = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(8.4), Inches(2.9))
    findings_frame = findings_text.text_frame
    findings_frame.text = "Critical Category-Specific Actions:"
    findings_frame.paragraphs[0].font.bold = True
    findings_frame.paragraphs[0].font.size = Pt(14)
    
    category_actions = [
        ("Customer Relations - Close Account (22.4 impact score):",
         ["Add context filters for retention discussions",
          "Distinguish between complaint and service request",
          "Review query complexity and reduce OR clauses"]),
        ("Credit Card ATM - Unclassified (14.6 impact score):",
         ["Emergency rule audit for 'unclassified' category",
          "Improve classification specificity",
          "Add equipment vs service distinction"]),
        ("Credit Card ATM - Rejected/Declined (12.5 impact score):",
         ["Separate technical issues from complaints",
          "Add merchant vs bank responsibility filters",
          "Implement category-specific negation rules"])
    ]
    
    for category, actions in category_actions:
        p = findings_frame.add_paragraph()
        p.text = f"\n{category}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_BLUE
        
        for action in actions:
            p = findings_frame.add_paragraph()
            p.text = f"  → {action}"
            p.font.size = Pt(11)
    
    # Expected impact
    impact_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.4))
    impact_frame = impact_text.text_frame
    impact_frame.text = "Expected Impact: +25-30% precision for these categories | Implementation: High priority (immediate resources required)"
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.size = Pt(13)
    impact_frame.paragraphs[0].font.color.rgb = GREEN
    
    # Slide 7: December 2024 Validation Investigation
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 7)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "December 2024 Validation Crisis - PROCESS CRITICAL"
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RED
    
    # Alert box
    alert_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.8)
    )
    alert_box.fill.solid()
    alert_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    alert_box.line.color.rgb = RED
    alert_box.line.width = Pt(3)
    
    alert_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.6))
    alert_frame = alert_text.text_frame
    alert_frame.text = "10.9% validation agreement drop (86.7% → 75.8%) affecting process credibility"
    alert_frame.paragraphs[0].font.size = Pt(16)
    alert_frame.paragraphs[0].font.bold = True
    
    # Validation trends table
    rows = 7
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.2), Inches(9), Inches(2.2)).table
    
    headers = ['Month', 'Agreement Rate', 'Sample Size', 'Status']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    
    validation_data = [
        ['October 2024', '86.7%', '83', 'Normal'],
        ['November 2024', '86.0%', '107', 'Normal'],
        ['December 2024', '75.8%', '128', 'PROBLEM'],
        ['January 2025', '85.0%', '127', 'Recovered'],
        ['February 2025', '82.4%', '136', 'Below Target'],
        ['March 2025', '85.7%', '133', 'Normal']
    ]
    
    for i, row_data in enumerate(validation_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = value
            if j == 3:  # Status column
                if value == 'PROBLEM':
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RED
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                elif value == 'Below Target':
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ORANGE
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Critical agreement issues
    issues_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.7),
        Inches(4.2), Inches(1.8)
    )
    issues_box.fill.solid()
    issues_box.fill.fore_color.rgb = RGBColor(255, 250, 250)
    
    issues_text = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(3.6), Inches(1.4))
    issues_frame = issues_text.text_frame
    issues_frame.text = "Categories with Critical Agreement:"
    issues_frame.paragraphs[0].font.bold = True
    issues_frame.paragraphs[0].font.size = Pt(12)
    
    critical_categories = [
        "• Customer Relations 'action not taken': 50.0%",
        "• EService 'login and registration': 58.3%",
        "• Payments 'missing precisely': 60.0%"
    ]
    
    for cat in critical_categories:
        p = issues_frame.add_paragraph()
        p.text = cat
        p.font.size = Pt(11)
        p.font.color.rgb = RED
    
    # Actions
    actions_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(4.7),
        Inches(4.5), Inches(1.8)
    )
    actions_box.fill.solid()
    actions_box.fill.fore_color.rgb = LIGHT_BLUE
    
    actions_text = slide.shapes.add_textbox(Inches(5.3), Inches(4.9), Inches(3.9), Inches(1.4))
    actions_frame = actions_text.text_frame
    actions_frame.text = "Immediate Actions:"
    actions_frame.paragraphs[0].font.bold = True
    actions_frame.paragraphs[0].font.size = Pt(12)
    
    validation_actions = [
        "• Interview December 2024 team",
        "• Recalibrate all validators",
        "• Implement consistency monitoring",
        "• Increase secondary validation to 30%"
    ]
    
    for action in validation_actions:
        p = actions_frame.add_paragraph()
        p.text = action
        p.font.size = Pt(11)
    
    # Slide 8: Volume-Performance Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 8)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Volume-Performance Anti-Correlation Management"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Key finding
    finding_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.6)
    )
    finding_box.fill.solid()
    finding_box.fill.fore_color.rgb = YELLOW_BG
    
    finding_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.4))
    finding_frame = finding_text.text_frame
    finding_frame.text = "Evidence: -0.135 correlation (high volume = lower precision)"
    finding_frame.paragraphs[0].font.size = Pt(14)
    finding_frame.paragraphs[0].font.bold = True
    
    # High-volume low-precision matrix
    rows = 6
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(2.5)).table
    
    headers = ['Category', 'Volume', 'Precision Gap', 'Impact Score']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    
    volume_data = [
        ['Customer Relations "close account"', '162', '13.8%', '22.4'],
        ['Credit Card ATM "unclassified"', '48', '30.4%', '14.6'],
        ['Credit Card ATM "rejected/declined"', '55', '22.7%', '12.5'],
        ['Fraud "general dissatisfaction"', '147', '6.7%', '9.9'],
        ['Credit Card ATM "travel notification"', '98', '9.8%', '9.6']
    ]
    
    for i, row_data in enumerate(volume_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = value
            if j == 3 and float(value) > 15:  # High impact score
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 235, 235)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Strategic actions
    strategy_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.8),
        Inches(9), Inches(1.6)
    )
    strategy_box.fill.solid()
    strategy_box.fill.fore_color.rgb = LIGHT_BLUE
    
    strategy_text = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(8.4), Inches(1.2))
    strategy_frame = strategy_text.text_frame
    strategy_frame.text = "Strategic Actions:"
    strategy_frame.paragraphs[0].font.bold = True
    strategy_frame.paragraphs[0].font.size = Pt(14)
    
    strategies = [
        "1. Implement volume-based precision thresholds",
        "2. Prioritize rule optimization by volume × precision gap score",
        "3. Create category-specific monitoring alerts",
        "4. Develop volume-based rule sensitivity adjustments"
    ]
    
    for strategy in strategies:
        p = strategy_frame.add_paragraph()
        p.text = strategy
        p.font.size = Pt(11)
    
    # Expected impact
    impact_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.4))
    impact_frame = impact_text.text_frame
    impact_frame.text = "Expected Impact: +10-15% precision for high-volume categories | ROI: Highest impact per effort ratio"
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.size = Pt(13)
    impact_frame.paragraphs[0].font.color.rgb = GREEN
    
    # Slide 9: Temporal Operational Patterns
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 9)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Temporal Operational Patterns"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Day of week patterns
    dow_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(4.2), Inches(2.3)
    )
    dow_box.fill.solid()
    dow_box.fill.fore_color.rgb = RGBColor(250, 250, 255)
    
    dow_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(3.6), Inches(1.9))
    dow_frame = dow_text.text_frame
    dow_frame.text = "Day-of-Week FP Rates:"
    dow_frame.paragraphs[0].font.bold = True
    dow_frame.paragraphs[0].font.size = Pt(13)
    
    dow_data = [
        "Monday: 27.1% (highest - 406 records)",
        "Tuesday: 26.1% (314 records)",
        "Wednesday: 23.7% (211 records)",
        "Thursday: 25.6% (550 records)",
        "Friday: 22.9% (lowest - 481 records)",
        "Weekend avg: 24.9% (893 records)"
    ]
    
    for day in dow_data:
        p = dow_frame.add_paragraph()
        p.text = f"• {day}"
        p.font.size = Pt(11)
        if "highest" in day:
            p.font.color.rgb = RED
        elif "lowest" in day:
            p.font.color.rgb = GREEN
    
    # Week of month patterns
    wom_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.2),
        Inches(4.5), Inches(2.3)
    )
    wom_box.fill.solid()
    wom_box.fill.fore_color.rgb = RGBColor(255, 250, 245)
    
    wom_text = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(3.9), Inches(1.9))
    wom_frame = wom_text.text_frame
    wom_frame.text = "Week-of-Month FP Rates:"
    wom_frame.paragraphs[0].font.bold = True
    wom_frame.paragraphs[0].font.size = Pt(13)
    
    wom_data = [
        "Week 1: 19.2% (443 records)",
        "Week 2: 25.8% (824 records)",
        "Week 3: 26.2% (778 records)",
        "Week 4: 26.9% (peak - 717 records)",
        "Week 5: 22.6% (93 records)"
    ]
    
    for week in wom_data:
        p = wom_frame.add_paragraph()
        p.text = f"• {week}"
        p.font.size = Pt(11)
        if "peak" in week:
            p.font.color.rgb = RED
    
    # Month-end discovery
    discovery_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.8),
        Inches(9), Inches(1.4)
    )
    discovery_box.fill.solid()
    discovery_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    discovery_box.line.color.rgb = GREEN
    discovery_box.line.width = Pt(2)
    
    discovery_text = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(8.4), Inches(1))
    discovery_frame = discovery_text.text_frame
    discovery_frame.text = "Unexpected Positive Finding - Month-End Performance:"
    discovery_frame.paragraphs[0].font.bold = True
    discovery_frame.paragraphs[0].font.size = Pt(13)
    discovery_frame.paragraphs[0].font.color.rgb = GREEN
    
    month_end_data = [
        "• Regular Days: 25.4% FP rate, 74.6% TP rate, 0.73 qualifying language",
        "• Month-End Days: 23.7% FP rate, 76.3% TP rate, 0.68 qualifying language",
        "• Month-end shows +1.9% precision improvement with less ambiguous language"
    ]
    
    for finding in month_end_data:
        p = discovery_frame.add_paragraph()
        p.text = finding
        p.font.size = Pt(11)
    
    # Actions
    actions_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.4),
        Inches(9), Inches(1.2)
    )
    actions_box.fill.solid()
    actions_box.fill.fore_color.rgb = LIGHT_BLUE
    
    actions_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(8.4), Inches(0.8))
    actions_frame = actions_text.text_frame
    actions_frame.text = "Strategic Actions: Investigate Monday differences | Apply month-end practices year-round | Implement day-specific thresholds"
    actions_frame.paragraphs[0].font.size = Pt(12)
    
    # Slide 10: Transcript Content Discrimination
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 10)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Transcript Content Discrimination - STATISTICAL SIGNIFICANCE"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Length analysis
    length_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(4.2), Inches(2)
    )
    length_box.fill.solid()
    length_box.fill.fore_color.rgb = RGBColor(255, 250, 250)
    
    length_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(3.6), Inches(1.6))
    length_frame = length_text.text_frame
    length_frame.text = "Transcript Length Analysis:"
    length_frame.paragraphs[0].font.bold = True
    length_frame.paragraphs[0].font.size = Pt(13)
    
    length_data = [
        "True Positives: 5,366 chars avg",
        "False Positives: 3,778 chars avg",
        "Difference: -1,588 chars (29.6% shorter)",
        "p-value: < 0.000001 (99.9999% confidence)"
    ]
    
    for item in length_data:
        p = length_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        if "p-value" in item:
            p.font.bold = True
            p.font.color.rgb = GREEN
    
    # Word ratio analysis
    ratio_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.2),
        Inches(4.5), Inches(2)
    )
    ratio_box.fill.solid()
    ratio_box.fill.fore_color.rgb = RGBColor(250, 250, 255)
    
    ratio_text = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(3.9), Inches(1.6))
    ratio_frame = ratio_text.text_frame
    ratio_frame.text = "Conversation Balance:"
    ratio_frame.paragraphs[0].font.bold = True
    ratio_frame.paragraphs[0].font.size = Pt(13)
    
    ratio_data = [
        "TP Customer Words: 554.42 avg",
        "FP Customer Words: 357.03 avg",
        "Customer-Agent Ratio:",
        "  TP: 0.93 | FP: 0.787 (-14.3% diff)"
    ]
    
    for item in ratio_data:
        p = ratio_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
    
    # Pattern risk assessment
    pattern_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.4),
        Inches(9), Inches(2)
    )
    pattern_box.fill.solid()
    pattern_box.fill.fore_color.rgb = RGBColor(255, 253, 240)
    
    pattern_text = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.6))
    pattern_frame = pattern_text.text_frame
    pattern_frame.text = "Advanced Pattern Risk Assessment (All patterns show risk factor <2.0):"
    pattern_frame.paragraphs[0].font.bold = True
    pattern_frame.paragraphs[0].font.size = Pt(13)
    
    patterns = [
        "• Politeness: 99.1% TP vs 98.8% FP (minimal discrimination)",
        "• Uncertainty: 69.5% TP vs 60.4% FP (moderate discrimination)",
        "• Frustration: 14.0% TP vs 5.4% FP (significant discrimination)",
        "• Hypotheticals: 45.5% TP vs 34.8% FP (moderate discrimination)"
    ]
    
    for pattern in patterns:
        p = pattern_frame.add_paragraph()
        p.text = pattern
        p.font.size = Pt(11)
    
    # Implementation strategy
    impl_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.6),
        Inches(9), Inches(1)
    )
    impl_box.fill.solid()
    impl_box.fill.fore_color.rgb = LIGHT_BLUE
    
    impl_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.6))
    impl_frame = impl_text.text_frame
    impl_frame.text = "Implementation: Set >2,500 char thresholds | Length-based confidence scoring | ML-based pattern detection"
    impl_frame.paragraphs[0].font.size = Pt(12)
    p = impl_frame.add_paragraph()
    p.text = "Expected Impact: +8-14% precision through content filtering and advanced patterns"
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    # Slide 11: Implementation Timeline & ROI
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_number(slide, 11)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Implementation Timeline & Expected ROI"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Timeline phases
    timeline_data = [
        ("Week 1-2", "CRITICAL FIXES", [
            "Deploy universal negation template",
            "Emergency category intervention",
            "December validation investigation"
        ], RED),
        ("Month 1", "SYSTEMATIC IMPROVEMENTS", [
            "Query optimization program",
            "Enhanced validation process",
            "Pattern-based improvements"
        ], ORANGE),
        ("Quarter 1", "STRATEGIC INITIATIVES", [
            "ML-based FP prediction",
            "Context-aware rule engine",
            "Semantic understanding layer"
        ], GREEN)
    ]
    
    y_pos = 1.3
    for period, phase, items, color in timeline_data:
        # Period marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(1.5), Inches(0.5)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        
        marker_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(1.5), Inches(0.5))
        marker_frame = marker_text.text_frame
        marker_frame.text = period
        marker_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        marker_frame.paragraphs[0].font.bold = True
        marker_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Phase content
        phase_text = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos), Inches(7.3), Inches(1.5))
        phase_frame = phase_text.text_frame
        phase_frame.text = phase
        phase_frame.paragraphs[0].font.bold = True
        phase_frame.paragraphs[0].font.size = Pt(12)
        phase_frame.paragraphs[0].font.color.rgb = color
        
        for item in items:
            p = phase_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
        
        y_pos += 1.8
    
    # ROI calculation
    roi_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.5),
        Inches(9), Inches(1.5)
    )
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    roi_box.line.color.rgb = GREEN
    roi_box.line.width = Pt(3)
    
    roi_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(8.4), Inches(1.1))
    roi_frame = roi_text.text_frame
    roi_frame.text = "Expected Outcomes & ROI:"
    roi_frame.paragraphs[0].font.bold = True
    roi_frame.paragraphs[0].font.size = Pt(14)
    roi_frame.paragraphs[0].font.color.rgb = GREEN
    
    roi_data = [
        "• Current Precision: 74.8% → Target: 85-89% (Expected Gain: +10-15%)",
        "• Annual Cost Savings: $2.4M (96.7% × $2.4M = $2.32M from context fixes alone)",
        "• Implementation: 3-5 FTEs for 12 weeks | Feasibility: YES - achieving target is feasible"
    ]
    
    for item in roi_data:
        p = roi_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
    
    # # Slide 12: Monitoring & Risk Management Framework
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    # add_slide_number(slide, 12)
    
    # title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    # title_frame = title_box.text_frame
    # title_frame.text = "Monitoring & Risk Management Framework"
    # title_frame.paragraphs[0].font.size = Pt(32)
    # title_frame.paragraphs[0].font.bold = True
    # title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Save presentation
    prs.save('Precision_Drop_Analysis_Presentation.pptx')
    print("Presentation created successfully: Precision_Drop_Analysis_Presentation.pptx")
    print(f"Total slides: {len(prs.slides)}")
    print("\nSlide contents:")
    print("1. Title Slide")
    print("2. Executive Summary")
    print("3. Performance vs. Targets")
    print("4. Critical Priority #1: Context-Insensitive Negation")
    print("5. Critical Priority #2: Agent Explanation Contamination")
    print("6. Emergency Category Review")
    print("7. Implementation Timeline")
    print("8. Monitoring & Risk Management")
    print("9. Key Recommendations & Next Steps")

# Create the presentation
if __name__ == "__main__":
    create_precision_analysis_presentation()
